import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
from openpyxl import Workbook
import streamlit as st
import plotly.express as px
import smtplib
from email.mime.text import MIMEText
from twilio.rest import Client as TwilioClient

class ReportGenerator:
    """
    PPTX 및 Excel 리포트 자동 생성
    """
    def __init__(self, template_pptx: str = None):
        self.template_pptx = template_pptx

    def generate_pptx(self, data: pd.DataFrame, output_path: str):
        """
        월간/분기별 PPT 보고서 생성
        :param data: 요약 데이터프레임 (KPI, 이상 수 등)
        :param output_path: 저장할 PPTX 경로
        """
        prs = Presentation(self.template_pptx) if self.template_pptx else Presentation()
        # 제목 슬라이드
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "월간 이상 탐지 리포트"
        # 테이블 삽입
        rows, cols = data.shape
        table = slide.shapes.add_table(rows+1, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)).table
        # 헤더
        for j, col in enumerate(data.columns):
            table.cell(0, j).text = col
        # 데이터
        for i, row in enumerate(data.itertuples(index=False), 1):
            for j, val in enumerate(row):
                table.cell(i, j).text = str(val)
        prs.save(output_path)

    def generate_excel(self, data: pd.DataFrame, output_path: str):
        """
        월간/분기별 Excel 보고서 생성
        :param data: 요약 데이터프레임
        :param output_path: 저장할 XLSX 경로
        """
        wb = Workbook()
        ws = wb.active
        ws.title = "Report"
        # 헤더
        ws.append(list(data.columns))
        # 데이터
        for row in data.itertuples(index=False):
            ws.append(list(row))
        wb.save(output_path)


def run_dashboard(data_csv: str):
    """
    Streamlit 대시보드 실행
    :param data_csv: 이상 탐지 결과 CSV 경로
    """
    st.title("실시간 이상 탐지 대시보드")
    df = pd.read_csv(data_csv, index_col=0)
    df['timestamp'] = pd.to_datetime(df['timestamp'])
    fig = px.line(df, x='timestamp', y='anomaly_count', title='시간별 이상 탐지 수')
    st.plotly_chart(fig)
    # KPI 필터
    metric = st.selectbox('Metric 선택', df.columns.drop(['timestamp','anomaly_count']))
    fig2 = px.histogram(df, x=metric, title=f'{metric} 분포')
    st.plotly_chart(fig2)

class AlertDispatcher:
    """
    다중 채널 알림 전송 (Slack, Email, SMS)
    """
    def __init__(self, slack_webhook: str = None, email_cfg: dict = None, sms_cfg: dict = None):
        self.slack_webhook = slack_webhook
        self.email_cfg = email_cfg
        self.sms_cfg = sms_cfg

    def send_slack(self, message: str):
        import requests
        if not self.slack_webhook:
            return
        payload = {"text": message}
        requests.post(self.slack_webhook, json=payload)

    def send_email(self, to_addrs: list, subject: str, body: str):
        cfg = self.email_cfg or {}
        msg = MIMEText(body)
        msg['Subject'] = subject
        msg['From'] = cfg.get('from_addr')
        msg['To'] = ','.join(to_addrs)
        with smtplib.SMTP(cfg.get('smtp_host'), cfg.get('smtp_port', 25)) as server:
            if cfg.get('starttls'):
                server.starttls()
            if cfg.get('username'):
                server.login(cfg.get('username'), cfg.get('password'))
            server.sendmail(msg['From'], to_addrs, msg.as_string())

    def send_sms(self, to_number: str, body: str):
        cfg = self.sms_cfg or {}
        client = TwilioClient(cfg.get('account_sid'), cfg.get('auth_token'))
        client.messages.create(
            body=body,
            from_=cfg.get('from_number'),
            to=to_number
        )

    def notify_all(self, message: str, email_to: list = None, sms_to: list = None):
        self.send_slack(message)
        if email_to:
            self.send_email(email_to, 'Risk Alert', message)
        if sms_to:
            for num in sms_to:
                self.send_sms(num, message)
